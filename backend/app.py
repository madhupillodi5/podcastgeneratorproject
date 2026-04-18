"""
╔══════════════════════════════════════════════════════════════════╗
║        AI PODCAST GENERATOR — BACKEND (Flask + VS Code)         ║
║  Stack: Flask · Groq · ElevenLabs · MongoDB · JWT · bcrypt      ║
║         pypdf · python-docx · python-pptx                       ║
║         pytesseract · pdf2image  (OCR — optional)               ║
║                                                                  ║
║  Run locally with:  python app.py                               ║
║  Accessible at:     http://localhost:5000                        ║
║                                                                  ║
║  Store API keys in a .env file (see .env.example)               ║
╚══════════════════════════════════════════════════════════════════╝
"""

import os, io, re, struct, base64, logging, textwrap, time, uuid
from datetime import datetime, timezone, timedelta
from concurrent.futures import ThreadPoolExecutor, as_completed
from functools import wraps
from pymongo.errors import DuplicateKeyError
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from dotenv import load_dotenv
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
from groq import Groq
from elevenlabs import ElevenLabs
import jwt
import bcrypt

# ── Load .env (API keys stored in .env file, never hard-coded) ────
load_dotenv()
# ── OCR imports (graceful fallback) ──────────────────────────────
try:
    import pytesseract
    from pdf2image import convert_from_bytes
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# ── MongoDB (optional — falls back to in-memory if unavailable) ──
MONGO_URI = os.environ.get("MONGO_URI", "")
db = None
users_col = None
podcasts_col = None

if MONGO_URI:
    try:
        from pymongo import MongoClient
        from pymongo.server_api import ServerApi
        client = MongoClient(MONGO_URI, server_api=ServerApi('1'), serverSelectionTimeoutMS=5000)
        client.admin.command('ping')
        db = client["podcastai"]
        users_col    = db["users"]
        podcasts_col = db["podcasts"]
        # Indexes
        users_col.create_index("email", unique=True)
        podcasts_col.create_index("user_id")
        podcasts_col.create_index("share_id", unique=True, sparse=True)
        print("✅  MongoDB connected.")
    except Exception as e:
        print(f"⚠️   MongoDB unavailable ({e}). Using in-memory storage.")
        db = None
else:
    print("⚠️   MONGO_URI not set. Using in-memory storage (data lost on restart).")

# ── In-memory fallback storage ────────────────────────────────────
_mem_users    = {}   # email -> user dict
_mem_podcasts = []   # list of podcast dicts

# ── Logging ──────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  [%(levelname)s]  %(message)s",
    datefmt="%H:%M:%S"
)
log = logging.getLogger("PodcastGen")

if OCR_AVAILABLE:
    log.info("OCR support: ENABLED")
else:
    log.warning("OCR support: DISABLED — install tesseract + pdf2image for scanned PDFs")

# ── Credentials from .env ─────────────────────────────────────────
def _env(key, required=True, default=""):
    val = os.environ.get(key, default)
    if required and not val:
        raise EnvironmentError(
            f"Missing required env variable: {key}\n"
            f"Add it to your .env file. See .env.example for reference."
        )
    return val

GROQ_API_KEY       = _env("GROQ_API_KEY")
ELEVENLABS_API_KEY = _env("ELEVENLABS_API_KEY")
JWT_SECRET         = _env("JWT_SECRET", default="dev-secret-change-in-production-32chars")
FLASK_SECRET_KEY   = _env("FLASK_SECRET_KEY", default="dev-flask-secret")

# ── AI Clients ────────────────────────────────────────────────────
groq_client       = Groq(api_key=GROQ_API_KEY)
elevenlabs_client = ElevenLabs(api_key=ELEVENLABS_API_KEY)

GROQ_MODEL   = "llama-3.3-70b-versatile"
VOICE_ALEX   = "nPczCjzI2devNBz1zQrb"   # Brian  – warm male
VOICE_JORDAN = "cgSgspJ2msm6clMCkdW9"   # Jessica – bright female

# ── Flask App ────────────────────────────────────────────────────
app = Flask(__name__, static_folder="../frontend", template_folder="../frontend")
app.secret_key = FLASK_SECRET_KEY

# Allow requests from localhost frontend (same origin in production)
CORS(app, supports_credentials=True, resources={r"/api/*": {"origins": "*"}})

ALLOWED_EXTENSIONS = {"pdf", "txt", "docx", "doc", "ppt", "pptx"}

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

# ════════════════════════════════════════════════════════════════
#   AUTH HELPERS
# ════════════════════════════════════════════════════════════════

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode("utf-8"), hashed.encode("utf-8"))

def create_token(user_id: str, email: str, name: str) -> str:
    now = datetime.now(timezone.utc)

    payload = {
        "sub": user_id,
        "email": email,
        "name": name,
        "iat": int(now.timestamp()),
        "exp": int((now + timedelta(days=30)).timestamp()),
    }

    return jwt.encode(payload, JWT_SECRET, algorithm="HS256")

def decode_token(token: str) -> dict:
    return jwt.decode(token, JWT_SECRET, algorithms=["HS256"])

def require_auth(f):
    """Decorator: extracts JWT from Authorization header."""
    @wraps(f)
    def decorated(*args, **kwargs):
        auth = request.headers.get("Authorization", "")
        if not auth.startswith("Bearer "):
            return jsonify({"error": "Authentication required."}), 401
        token = auth.split(" ", 1)[1]
        try:
            payload = decode_token(token)
            request.user = payload
        except jwt.ExpiredSignatureError:
            return jsonify({"error": "Session expired. Please log in again."}), 401
        except jwt.InvalidTokenError:
            return jsonify({"error": "Invalid token."}), 401
        return f(*args, **kwargs)
    return decorated

# ── DB helpers (abstract over Mongo vs in-memory) ────────────────

def db_find_user(email: str):
    if db is not None:
        return users_col.find_one({"email": email})
    return _mem_users.get(email)

def db_create_user(email: str, name: str, password: str) -> str:
    uid = str(uuid.uuid4())
    user = {
        "_id": uid,
        "email": email,
        "name": name,
        "password": hash_password(password),
        "created_at": datetime.utcnow().isoformat()
    }
    if db is not None:
        try:
            users_col.insert_one(user)
        except DuplicateKeyError:
            raise ValueError("An account with this email already exists.")
    else:
        _mem_users[email] = user
    return uid

def db_get_user_podcasts(user_id: str) -> list:
    if db is not None:
        docs = list(podcasts_col.find({"user_id": user_id}, {"_id": 0}).sort("created_at", -1).limit(100))
        return docs
    return [p for p in _mem_podcasts if p.get("user_id") == user_id]

def db_save_podcast(entry: dict):
    if db is not None:
        podcasts_col.insert_one(entry)
    else:
        _mem_podcasts.insert(0, entry)

def db_delete_podcast(podcast_id: str, user_id: str):
    if db is not None:
        podcasts_col.delete_one({"podcast_id": podcast_id, "user_id": user_id})
    else:
        global _mem_podcasts
        _mem_podcasts = [p for p in _mem_podcasts if not (p["podcast_id"] == podcast_id and p["user_id"] == user_id)]

def db_find_by_share_id(share_id: str):
    if db is not None:
        return podcasts_col.find_one({"share_id": share_id}, {"_id": 0})
    return next((p for p in _mem_podcasts if p.get("share_id") == share_id), None)

def db_rename_podcast(podcast_id: str, user_id: str, new_title: str):
    if db is not None:
        podcasts_col.update_one({"podcast_id": podcast_id, "user_id": user_id}, {"$set": {"title": new_title}})
    else:
        for p in _mem_podcasts:
            if p["podcast_id"] == podcast_id and p["user_id"] == user_id:
                p["title"] = new_title
                break

# ════════════════════════════════════════════════════════════════
#   AUTH ROUTES
# ════════════════════════════════════════════════════════════════

@app.route("/api/auth/signup", methods=["POST"])
def signup():
    log.info("DB object is: %s", db)
    log.info("users_col is: %s", users_col)
    data = request.get_json() or {}
    name     = (data.get("name") or "").strip()
    email    = (data.get("email") or "").strip().lower()
    password = (data.get("password") or "")

    if not name or not email or not password:
        return jsonify({"error": "Name, email and password are required."}), 400
    if len(password) < 8:
        return jsonify({"error": "Password must be at least 8 characters."}), 400
    if db_find_user(email):
        return jsonify({"error": "An account with this email already exists."}), 409

    try:
        uid = db_create_user(email, name, password)
    except ValueError as e:
        return jsonify({"error": str(e)}), 409
    except Exception as e:
        log.exception("Signup DB error: %s", e)
        return jsonify({"error": "Failed to create account. Please try again."}), 500

    token = create_token(uid, email, name)
    return jsonify({"token": token, "user": {"id": uid, "email": email, "name": name}}), 201


@app.route("/api/auth/login", methods=["POST"])
def login():
    data = request.get_json() or {}
    email    = (data.get("email") or "").strip().lower()
    password = (data.get("password") or "")

    user = db_find_user(email)
    if not user or not verify_password(password, user["password"]):
        return jsonify({"error": "Invalid email or password."}), 401

    uid   = str(user["_id"])
    token = create_token(uid, email, user["name"])
    return jsonify({"token": token, "user": {"id": uid, "email": email, "name": user["name"]}}), 200


@app.route("/api/auth/me", methods=["GET"])
@require_auth
def me():
    return jsonify({"user": {"id": request.user["sub"], "email": request.user["email"], "name": request.user["name"]}}), 200

# ════════════════════════════════════════════════════════════════
#   PODCAST HISTORY ROUTES
# ════════════════════════════════════════════════════════════════

@app.route("/api/podcasts", methods=["GET"])
@require_auth
def get_podcasts():
    items = db_get_user_podcasts(request.user["sub"])
    # Don't send full audio in list — only metadata + script preview
    result = []
    for p in items:
        result.append({
            "podcast_id": p.get("podcast_id"),
            "title":      p.get("title") or p.get("filename", "Untitled"),
            "filename":   p.get("filename"),
            "created_at": p.get("created_at"),
            "share_id":   p.get("share_id"),
            "script_preview": (p.get("script") or "")[:200],
        })
    return jsonify({"podcasts": result}), 200


@app.route("/api/podcasts/<podcast_id>", methods=["GET"])
@require_auth
def get_podcast(podcast_id):
    items = db_get_user_podcasts(request.user["sub"])
    p = next((x for x in items if x.get("podcast_id") == podcast_id), None)
    if not p:
        return jsonify({"error": "Not found."}), 404
    return jsonify({"podcast": {k: v for k, v in p.items() if k != "_id"}}), 200


@app.route("/api/podcasts/<podcast_id>", methods=["DELETE"])
@require_auth
def delete_podcast(podcast_id):
    db_delete_podcast(podcast_id, request.user["sub"])
    return jsonify({"ok": True}), 200


@app.route("/api/podcasts/<podcast_id>/rename", methods=["PATCH"])
@require_auth
def rename_podcast(podcast_id):
    data = request.get_json() or {}
    new_title = (data.get("title") or "").strip()
    if not new_title:
        return jsonify({"error": "Title is required."}), 400
    db_rename_podcast(podcast_id, request.user["sub"], new_title)
    return jsonify({"ok": True}), 200


@app.route("/api/share/<share_id>", methods=["GET"])
def get_shared(share_id):
    """Public endpoint — no auth required."""
    p = db_find_by_share_id(share_id)
    if not p:
        return jsonify({"error": "Shared podcast not found."}), 404
    return jsonify({
        "title":    p.get("title") or p.get("filename", "Podcast"),
        "script":   p.get("script", ""),
        "audio":    p.get("audio", ""),
        "filename": p.get("filename", ""),
        "created_at": p.get("created_at", ""),
    }), 200

# ════════════════════════════════════════════════════════════════
#   TEXT EXTRACTION
# ════════════════════════════════════════════════════════════════

def extract_text(file_bytes, filename):
    ext = filename.rsplit(".", 1)[1].lower()
    log.info("Extracting text from .%s (%.1f KB)...", ext, len(file_bytes) / 1024)

    if ext == "pdf":
        text = _extract_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        doc  = DocxDocument(io.BytesIO(file_bytes))
        text = "\n".join(p.text for p in doc.paragraphs)
    elif ext in ("ppt", "pptx"):
        prs  = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        lines.append(para.text)
        text = "\n".join(lines)
    elif ext == "txt":
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: .{ext}")

    cleaned = _clean_text(text)
    log.info("Extracted %d chars / ~%d words.", len(cleaned), len(cleaned.split()))

    if not cleaned:
        raise ValueError(
            "No text could be extracted from the file. "
            "If it is a scanned PDF, install tesseract and pdf2image."
        )
    return cleaned


def _extract_pdf(file_bytes):
    reader      = pypdf.PdfReader(io.BytesIO(file_bytes))
    total_pages = len(reader.pages)
    log.info("PDF has %d page(s).", total_pages)

    pages_text, empty_pages = [], 0
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if not page_text.strip():
            empty_pages += 1
        else:
            pages_text.append(page_text)

    if total_pages > 0 and empty_pages / total_pages <= 0.5:
        return "\n".join(pages_text)

    if not OCR_AVAILABLE:
        raise ValueError(
            "This PDF appears to be scanned (image-only). "
            "Install tesseract-ocr and pdf2image to enable OCR support."
        )
    return _ocr_pdf(file_bytes, total_pages)


def _ocr_pdf(file_bytes, total_pages):
    dpi = 150 if total_pages > 30 else 200
    log.info("OCR: %d pages at %d DPI...", total_pages, dpi)
    t0 = time.time()
    images = convert_from_bytes(file_bytes, dpi=dpi, fmt="jpeg", thread_count=4)
    results = {}

    def _ocr_page(args):
        idx, img = args
        try:
            return idx, pytesseract.image_to_string(img, lang="eng", config="--oem 3 --psm 3")
        except Exception as e:
            return idx, ""

    with ThreadPoolExecutor(max_workers=4) as ex:
        for idx, text in ex.map(_ocr_page, enumerate(images)):
            results[idx] = text

    combined = "\n".join(results.get(i, "") for i in range(len(images)))
    log.info("OCR done: %d chars in %.1fs.", len(combined), time.time() - t0)
    if not combined.strip():
        raise ValueError("OCR extracted no text. The PDF may be very low quality.")
    return combined


def _clean_text(raw):
    text = re.sub(r"[ \t]+", " ", raw)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[^\x20-\x7E\n]", "", text)
    return text.strip()

# ════════════════════════════════════════════════════════════════
#   BOILERPLATE FILTERING
# ════════════════════════════════════════════════════════════════

_BOILERPLATE_LINE_PATTERNS = [
    r"^page\s*(no|number)?[\s.:]*\d*$", r"^-\s*\d+\s*-$",
    r"^\d+\s*\|?\s*p\s*a\s*g\s*e$",
    r"signature\s*of\s*(the\s*)?(faculty|examiner|supervisor|guide|hod)",
    r"^(faculty|examiner|hod|supervisor|guide)\s*(signature)?[\s.:_-]*$",
    r"verified\s*by", r"approved\s*by",
    r"roll\s*(no|number)\s*[:\-.]?\s*\d*",
    r"register\s*(no|number)\s*[:\-.]?\s*\d*",
    r"^(table\s*of\s*contents?|index|contents?)$",
    r"(university|institute|college|department)\s*of\s*(technology|science|engineering|arts)",
    r"(autonomous|affiliated\s*to|accredited\s*by)",
    r"(naac|nba|ugc|aicte)\s*(accredited|approved|recognized)?",
    r"(academic\s*year|batch)\s*[:\-.]?\s*\d{4}",
    r"^about\s*(the\s*)?(author|writer)s?$",
    r"all\s*rights?\s*reserved", r"copyright\s*\©?\s*\d{4}",
    r"isbn\s*[:\-.]?\s*[\d\-]+",
    r"www\.[a-z0-9\-]+\.[a-z]{2,}", r"^https?://",
    r"^[_\-=*#~.]{3,}$", r"^[\W\s]{0,3}$",
]
_BOILERPLATE_LINE_RE = re.compile("|".join(_BOILERPLATE_LINE_PATTERNS), re.IGNORECASE)
_BOILERPLATE_PARA_RE = re.compile(
    r"(table\s*of\s*contents?|pin\s*code|phone\s*no|fax\s*no)", re.IGNORECASE
)

def _regex_filter(text):
    paragraphs = re.split(r"\n{2,}", text)
    kept = []
    for para in paragraphs:
        if _BOILERPLATE_PARA_RE.search(para):
            continue
        lines = para.splitlines()
        good = [l for l in lines if not _BOILERPLATE_LINE_RE.match(l.strip())]
        if len(good) < 2 and len(lines) > 3:
            continue
        cleaned = "\n".join(good).strip()
        if cleaned:
            kept.append(cleaned)
    return "\n\n".join(kept)

LLM_FILTER_PROMPT = """\
Is this text chunk MEANINGFUL CONTENT (concepts, data, arguments, findings) worth discussing in a podcast?
Or is it BOILERPLATE (table of contents, author bio, cover page, signatures, page numbers, copyright notices)?
Reply ONLY with: KEEP or DISCARD

CHUNK:
{chunk}
"""

def _llm_filter_chunks(chunks):
    kept = []
    for chunk in chunks:
        if len(chunk.split()) < 20:
            continue
        try:
            v = _groq(LLM_FILTER_PROMPT.format(chunk=chunk[:1500]), max_tokens=5).strip().upper()
            if v.startswith("KEEP"):
                kept.append(chunk)
        except Exception:
            kept.append(chunk)
    return kept

def filter_boilerplate(text):
    after_regex = _regex_filter(text)
    if not after_regex.strip():
        after_regex = text
    candidates = [c.strip() for c in textwrap.wrap(after_regex, width=1500, break_long_words=False) if c.strip()]
    kept = _llm_filter_chunks(candidates)
    return "\n\n".join(kept) if kept else after_regex

# ════════════════════════════════════════════════════════════════
#   SUMMARISATION & SCRIPT
# ════════════════════════════════════════════════════════════════

def chunk_text(text, size=2000):
    chunks = [c.strip() for c in textwrap.wrap(text, width=size, break_long_words=False) if c.strip()]
    log.info("Chunked into %d pieces.", len(chunks))
    return chunks

def _groq(prompt, max_tokens=1024, retries=3):
    for attempt in range(1, retries + 1):
        try:
            r = groq_client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=max_tokens,
                temperature=0.7,
            )
            return r.choices[0].message.content.strip()
        except Exception as e:
            log.warning("Groq attempt %d/%d: %s", attempt, retries, e)
            if attempt < retries:
                time.sleep(2 * attempt)
    raise RuntimeError("Groq API failed after all retries.")

def summarise_chunks(chunks):
    summaries = []
    for i, chunk in enumerate(chunks, 1):
        s = _groq(f"Summarise in 100-150 words, preserving key ideas:\n\n{chunk}")
        summaries.append(s)
        if i % 5 == 0:
            log.info("  Summarised %d/%d chunks.", i, len(chunks))
    return summaries

def hierarchical_summarise(summaries):
    if len(summaries) <= 30:
        return summaries
    batches = [summaries[i:i+7] for i in range(0, len(summaries), 7)]
    return [_groq("Combine into ONE paragraph of 80-100 words:\n\n" + "\n\n".join(f"- {s}" for s in b)) for b in batches]

SCRIPT_PROMPT = """\
Write a natural podcast conversation between Alex and Jordan using these summaries.

RULES:
- Every line: Alex: <text>   OR   Jordan: <text>
- No narration or stage directions
- Short, conversational sentences
- ~{target} words total
- Cover all key points

SUMMARIES:
{summaries}

BEGIN:
"""

def generate_script(summaries):
    combined = "\n\n".join(summaries)
    target   = max(300, min(sum(len(s) for s in summaries) // 20, 2400))
    if len(combined) > 40000:
        combined = combined[:40000] + "\n[truncated]"
    log.info("Generating script (~%d words)...", target)
    r = groq_client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[{"role": "user", "content": SCRIPT_PROMPT.format(summaries=combined, target=target)}],
        max_tokens=min(int(target * 1.4) + 400, 8192),
        temperature=0.8,
    )
    return r.choices[0].message.content.strip()

LINE_RE = re.compile(r"^(Alex|Jordan)\s*:\s*(.+)$", re.IGNORECASE)

def parse_script(script):
    dialogue, last = [], "Alex"
    for line in script.splitlines():
        line = line.strip()
        if not line:
            continue
        m = LINE_RE.match(line)
        speaker = m.group(1).capitalize() if m else last
        text    = m.group(2).strip() if m else line
        if not text:
            continue
        dialogue.append({"speaker": speaker, "text": text, "voice_id": VOICE_ALEX if speaker == "Alex" else VOICE_JORDAN})
        last = speaker
    return dialogue

# ════════════════════════════════════════════════════════════════
#   AUDIO GENERATION
# ════════════════════════════════════════════════════════════════

SAMPLE_RATE = 22050

def _split_tts(text, size=250):
    if len(text) <= size:
        return [text]
    sentences, chunks, cur = re.split(r"(?<=[.!?])\s+", text), [], ""
    for s in sentences:
        if len(cur) + len(s) + 1 <= size:
            cur = (cur + " " + s).strip()
        else:
            if cur:
                chunks.append(cur)
            cur = s[:size]
    if cur:
        chunks.append(cur)
    return chunks or [text[:size]]

def _tts(text, voice_id):
    return b"".join(elevenlabs_client.text_to_speech.convert(
        text=text, voice_id=voice_id, model_id="eleven_turbo_v2", output_format="pcm_22050"
    ))

def generate_audio(dialogue):
    segments = []
    pause = b"\x00\x00" * int(SAMPLE_RATE * 0.4)
    for i, entry in enumerate(dialogue, 1):
        for sub in _split_tts(entry["text"]):
            try:
                segments.append(_tts(sub, entry["voice_id"]))
            except Exception as e:
                log.warning("TTS line %d: %s", i, e)
        segments.append(pause)
    pcm = b"".join(segments)
    # Build WAV header
    data_size   = len(pcm)
    byte_rate   = SAMPLE_RATE * 2
    header = struct.pack("<4sI4s4sIHHIIHH4sI",
        b"RIFF", 36 + data_size, b"WAVE",
        b"fmt ", 16, 1, 1, SAMPLE_RATE, byte_rate, 2, 16,
        b"data", data_size)
    duration = data_size / byte_rate
    log.info("Audio: %.1f seconds (%.1f min).", duration, duration / 60)
    return header + pcm

# ════════════════════════════════════════════════════════════════
#   MAIN PIPELINE
# ════════════════════════════════════════════════════════════════

def run_pipeline(file_bytes, filename):
    t0 = time.time()
    log.info("=== Pipeline START: %s ===", filename)
    text          = extract_text(file_bytes, filename)
    filtered      = filter_boilerplate(text)
    chunks        = chunk_text(filtered)
    if not chunks:
        raise ValueError("No meaningful content found after filtering.")
    summaries     = summarise_chunks(chunks)
    final_sums    = hierarchical_summarise(summaries)
    script        = generate_script(final_sums)
    dialogue      = parse_script(script)
    if not dialogue:
        raise ValueError("Script parsing produced no dialogue.")
    wav           = generate_audio(dialogue)
    audio_b64     = base64.b64encode(wav).decode()
    log.info("=== Pipeline DONE in %.1fs ===", time.time() - t0)
    return {"script": script, "audio": audio_b64}

# ════════════════════════════════════════════════════════════════
#   MAIN PODCAST ROUTE
# ════════════════════════════════════════════════════════════════

@app.route("/api/generate-podcast", methods=["POST"])
@require_auth
def generate_podcast():
    if "file" not in request.files:
        return jsonify({"error": "No file in request."}), 400
    f = request.files["file"]
    if not f.filename or not allowed_file(f.filename):
        return jsonify({"error": "Unsupported file type. Use PDF, TXT, DOCX, or PPTX."}), 400
    file_bytes = f.read()
    if len(file_bytes) > 10 * 1024 * 1024:
        return jsonify({"error": "File exceeds 10 MB limit."}), 413
    if not file_bytes:
        return jsonify({"error": "Empty file."}), 400

    log.info("File: %s (%.1f KB) — user: %s", f.filename, len(file_bytes)/1024, request.user["email"])

    try:
        result = run_pipeline(file_bytes, f.filename)
    except ValueError as e:
        return jsonify({"error": str(e)}), 422
    except RuntimeError as e:
        return jsonify({"error": str(e)}), 502
    except Exception as e:
        log.exception("Pipeline error: %s", e)
        return jsonify({"error": "Internal server error."}), 500

    # Save to DB
    share_id   = str(uuid.uuid4()).replace("-", "")[:16]
    podcast_id = str(uuid.uuid4())
    entry = {
        "podcast_id":  podcast_id,
        "user_id":     request.user["sub"],
        "filename":    f.filename,
        "title":       f.filename.rsplit(".", 1)[0],
        "script":      result["script"],
        "audio":       result["audio"],
        "created_at":  datetime.utcnow().isoformat(),
        "share_id":    share_id,
    }
    db_save_podcast(entry)
    log.info("Saved podcast %s (share: %s)", podcast_id, share_id)

    return jsonify({
        "podcast_id": podcast_id,
        "script":     result["script"],
        "audio":      result["audio"],
        "share_id":   share_id,
    }), 200

# ════════════════════════════════════════════════════════════════
#   STATUS & FRONTEND SERVING
# ════════════════════════════════════════════════════════════════

@app.route("/api/status", methods=["GET"])
def status():
    return jsonify({
        "status":      "online",
        "service":     "AI Podcast Generator",
        "ocr_support": OCR_AVAILABLE,
        "db":          "mongodb" if db is not None else "memory",
    })

@app.route("/", defaults={"path": ""})
@app.route("/<path:path>")
def serve_frontend(path):
    """Serve the SPA for all non-API routes."""
    if path and path.startswith("api/"):
        return jsonify({"error": "Not found"}), 404
    return send_from_directory("../frontend", "index.html")

# ════════════════════════════════════════════════════════════════
#   ENTRY POINT
# ════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    print("\n" + "=" * 55)
    print("  🎙  AI PODCAST GENERATOR")
    print(f"  🌐  http://localhost:{port}")
    print(f"  🗄   DB: {'MongoDB Atlas' if db is not None else 'In-memory (set MONGO_URI for persistence)'}")
    print(f"  🔍  OCR: {'Enabled' if OCR_AVAILABLE else 'Disabled (install tesseract)'}")
    print("  ✅  API keys loaded from .env")
    print("=" * 55 + "\n")
    app.run(host="0.0.0.0", port=port, debug=os.environ.get("FLASK_ENV") == "development")
